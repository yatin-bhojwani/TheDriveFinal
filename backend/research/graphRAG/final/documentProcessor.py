import logging
import json
import re
import os
from datetime import datetime
from typing import Set, Optional, List, Dict, Tuple, Any
from dataclasses import dataclass
from enum import Enum
import PyPDF2
import docx
import openpyxl
import csv
from pptx import Presentation
from PIL import Image
import pytesseract
from pathlib import Path

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class SearchMode(Enum):
    """Types of GraphRAG search modes"""
    GLOBAL = "global"
    LOCAL = "local" 
    HYBRID = "hybrid"


class CommunityLevel(Enum):
    """Hierarchical community levels"""
    ENTITIES = 0      # Individual entities
    TOPICS = 1        # Small topic clusters (5-15 entities)
    THEMES = 2        # Mid-level themes (15-30 entities)
    DOMAINS = 3       # High-level domains (30+ entities)


@dataclass
class GraphRAGConfig:
    """Enhanced configuration for GraphRAG operations"""
    # Clustering parameters
    max_community_size: int = 100
    min_community_size: int = 3
    clustering_resolution: List[float] = None
    max_clustering_iterations: int = 50
    modularity_threshold: float = 0.3
    
    # Search parameters
    max_search_hops: int = 2
    global_search_timeout: int = 10
    local_search_timeout: int = 5
    hybrid_search_timeout: int = 8
    max_entities_per_search: int = 50
    
    # Context parameters
    max_context_tokens: int = 4000
    community_summary_length: int = 300
    max_communities_per_query: int = 15
    entity_description_length: int = 200
    
    # Performance parameters
    max_concurrent_searches: int = 5
    cache_ttl: int = 3600
    result_batch_size: int = 20
    
    # LLM parameters
    llm_model: str = "gemini-2.0-flash"
    embedding_model: str = "text-embedding-004"
    temperature: float = 0.1
    
    # GraphRAG specific parameters
    enable_community_detection: bool = True
    enable_hierarchical_search: bool = True
    entity_extraction_threshold: float = 0.7
    relationship_confidence_threshold: float = 0.6
    
    def __post_init__(self):
        if self.clustering_resolution is None:
            self.clustering_resolution = [0.1, 0.3, 0.5, 1.0, 2.0]


@dataclass
class Community:
    """Enhanced community representation"""
    community_id: str
    level: CommunityLevel
    entities: List[str]
    relationships: List[Dict]
    summary: str
    size: int
    modularity: float
    centrality_scores: Dict[str, float]
    topic_keywords: List[str]
    parent_community: Optional[str] = None
    child_communities: List[str] = None
    folder_id: str = ""
    
    def __post_init__(self):
        if self.child_communities is None:
            self.child_communities = []


@dataclass
class SearchResult:
    """Enhanced search operation result"""
    mode: SearchMode
    query: str
    folder_id: str
    entities: List[Dict]
    communities: List[Community]
    relationships: List[Dict]
    context_summary: str
    execution_time: float
    total_tokens: int
    relevance_scores: Dict[str, float]
    confidence_breakdown: Dict[str, float]
    source_attribution: List[Dict[str, str]]


class DocumentType(Enum):
    PDF = "pdf"
    DOCX = "docx"
    TXT = "txt"
    MD = "md"
    JSON = "json"
    CSV = "csv"
    XLSX = "xlsx"
    PPTX = "pptx"
    JPG = "jpg"
    PNG = "png"


@dataclass
class FolderScope:
    """Enhanced metadata for folder scope"""
    folder_path: str
    folder_id: str
    allowed_extensions: Set[str]
    max_file_size: int = 50 * 1024 * 1024  # 50MB 
    recursive_scan: bool = True
    # GraphRAG specific settings
    enable_entity_linking: bool = True
    content_deduplication: bool = True
    preserve_document_structure: bool = True


@dataclass
class DocumentMetadata:
    """Enhanced metadata for documents with GraphRAG features"""
    file_path: str
    file_name: str
    folder_path: str
    folder_id: str
    document_hash: str
    file_size: int
    modification_time: float
    document_type: DocumentType
    creation_timestamp: datetime
    processing_status: str = "pending"
    
    # GraphRAG specific metadata
    estimated_entity_count: int = 0
    content_complexity_score: float = 0.0
    structural_importance: float = 1.0
    topic_classification: str = "general"
    language_detected: str = "en"
    chunk_count: int = 0


class ContentAnalyzer:
    """Analyzes content for better GraphRAG processing"""
    
    def __init__(self):
        self.entity_patterns = [
            r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b',  # Proper nouns
            r'\b[A-Z]{2,}\b',  # Acronyms
            r'\b\w+(?:_\w+)+\b',  # Technical terms with underscores
        ]
        
        self.topic_keywords = {
            'mathematical': ['theorem', 'proof', 'equation', 'formula', 'algorithm', 'function'],
            'technical': ['system', 'architecture', 'implementation', 'framework', 'api'],
            'scientific': ['research', 'study', 'analysis', 'experiment', 'hypothesis'],
            'business': ['strategy', 'market', 'revenue', 'customer', 'product'],
            'legal': ['law', 'regulation', 'compliance', 'contract', 'policy']
        }
    
    def analyze_content(self, content: str) -> Dict[str, Any]:
        """Comprehensive content analysis for GraphRAG"""
        
        analysis = {
            'estimated_entities': self._count_potential_entities(content),
            'complexity_score': self._calculate_complexity(content),
            'topic_classification': self._classify_topic(content),
            'structure_analysis': self._analyze_structure(content),
            'keyword_density': self._calculate_keyword_density(content),
            'readability_score': self._calculate_readability(content)
        }
        
        return analysis
    
    def _count_potential_entities(self, content: str) -> int:
        """Count potential entities using pattern matching"""
        total_entities = 0
        
        for pattern in self.entity_patterns:
            matches = re.findall(pattern, content)
            total_entities += len(set(matches))  # Unique matches only
        
        return min(total_entities, len(content.split()) // 8)  # Cap at reasonable ratio
    
    def _calculate_complexity(self, content: str) -> float:
        """Calculate content complexity for processing priority"""
        
        factors = {
            'sentence_length': len(content.split('.')) / max(len(content.split()), 1),
            'vocabulary_diversity': len(set(content.lower().split())) / max(len(content.split()), 1),
            'technical_density': sum(1 for word in content.split() if len(word) > 8) / max(len(content.split()), 1)
        }
        
        return sum(factors.values()) / len(factors)
    
    def _classify_topic(self, content: str) -> str:
        """Classify content topic for better organization"""
        
        content_lower = content.lower()
        scores = {}
        
        for topic, keywords in self.topic_keywords.items():
            score = sum(1 for keyword in keywords if keyword in content_lower)
            if score > 0:
                scores[topic] = score
        
        if scores:
            return max(scores, key=scores.get)
        return 'general'
    
    def _analyze_structure(self, content: str) -> Dict[str, Any]:
        """Analyze document structure"""
        
        return {
            'paragraph_count': len([p for p in content.split('\n\n') if p.strip()]),
            'sentence_count': len([s for s in content.split('.') if s.strip()]),
            'has_headers': bool(re.search(r'^#+\s', content, re.MULTILINE)),
            'has_lists': bool(re.search(r'^\s*[-*•]\s', content, re.MULTILINE)),
            'has_code': bool(re.search(r'```|`[^`]+`', content)),
            'has_references': bool(re.search(r'\[\d+\]|\([^)]*\d{4}[^)]*\)', content))
        }
    
    def _calculate_keyword_density(self, content: str) -> Dict[str, float]:
        """Calculate density of important keywords"""
        
        words = content.lower().split()
        total_words = len(words)
        
        if total_words == 0:
            return {}
        
        # Count important word categories
        categories = {
            'technical': ['system', 'method', 'approach', 'framework', 'model'],
            'analytical': ['analysis', 'result', 'conclusion', 'finding', 'evidence'],
            'descriptive': ['description', 'overview', 'summary', 'introduction'],
            'relational': ['relationship', 'connection', 'correlation', 'association']
        }
        
        densities = {}
        for category, keywords in categories.items():
            count = sum(1 for word in words if word in keywords)
            densities[category] = count / total_words
        
        return densities
    
    def _calculate_readability(self, content: str) -> float:
        """Simple readability score (0-1, higher = more readable)"""
        
        if not content.strip():
            return 0.0
        
        sentences = [s.strip() for s in content.split('.') if s.strip()]
        words = content.split()
        
        if not sentences or not words:
            return 0.0
        
        avg_sentence_length = len(words) / len(sentences)
        # Normalize to 0-1 scale (assuming 15 words per sentence is optimal)
        readability = max(0, 1 - abs(avg_sentence_length - 15) / 20)
        
        return readability


class EnhancedDocumentProcessor:
    """Enhanced document processor with GraphRAG-specific features"""
    
    def __init__(self):
        self.supported_types = {
            '.pdf': DocumentType.PDF,
            '.docx': DocumentType.DOCX,
            '.doc': DocumentType.DOCX,
            '.txt': DocumentType.TXT,
            '.md': DocumentType.MD,
            '.json': DocumentType.JSON,
            '.csv': DocumentType.CSV,
            '.xlsx': DocumentType.XLSX,
            '.pptx': DocumentType.PPTX,
            '.jpg': DocumentType.JPG,
            '.jpeg': DocumentType.JPG,
            '.png': DocumentType.PNG
        }
        self.content_analyzer = ContentAnalyzer()
    
    def extract_content_with_analysis(self, file_path: str, doc_type: DocumentType) -> Tuple[str, Dict[str, Any]]:
        """Extract content and perform GraphRAG analysis"""
        
        # Extract base content
        content = self.extract_content(file_path, doc_type)
        
        if not content.strip():
            return content, {'error': 'No content extracted'}
        
        # Perform content analysis
        analysis = self.content_analyzer.analyze_content(content)
        
        # Add file-specific metadata
        analysis.update({
            'file_size': Path(file_path).stat().st_size,
            'extraction_method': doc_type.value,
            'content_length': len(content),
            'word_count': len(content.split())
        })
        
        return content, analysis
    
    def extract_content(self, file_path: str, doc_type: DocumentType) -> str:
        """Extract text content from document with enhanced error handling"""
        try:
            if doc_type == DocumentType.PDF:
                return self._extract_pdf_enhanced(file_path)
            elif doc_type == DocumentType.DOCX:
                return self._extract_docx_enhanced(file_path)
            elif doc_type == DocumentType.TXT:
                return self._extract_text_enhanced(file_path)
            elif doc_type == DocumentType.MD:
                return self._extract_markdown_enhanced(file_path)
            elif doc_type == DocumentType.JSON:
                return self._extract_json_enhanced(file_path)
            elif doc_type == DocumentType.CSV:
                return self._extract_csv_enhanced(file_path)
            elif doc_type == DocumentType.XLSX:
                return self._extract_xlsx_enhanced(file_path)
            elif doc_type == DocumentType.PPTX:
                return self._extract_pptx_enhanced(file_path)
            elif doc_type in [DocumentType.JPG, DocumentType.PNG]:
                return self._extract_image_enhanced(file_path)
            else:
                raise ValueError(f"Unsupported document type: {doc_type}")
        except Exception as e:
            return ""
    
    def _extract_pdf_enhanced(self, file_path: str) -> str:
        """Enhanced PDF extraction with better error handling and fallback methods"""
        content_parts = []
        
        try:
            import PyPDF2
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Check if PDF is encrypted
                if pdf_reader.is_encrypted:
                    logger.warning(f"PDF {file_path} is encrypted, attempting to decrypt")
                    try:
                        pdf_reader.decrypt('')  # Try empty password
                    except:
                        logger.error(f"Could not decrypt PDF {file_path}")
                        return ""
                
                # Extract metadata if available
                if pdf_reader.metadata:
                    title = pdf_reader.metadata.get('/Title', '')
                    if title:
                        content_parts.append(f"Document Title: {title}")
                
                # Extract page content with better error handling
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    try:
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            # Clean up the text
                            cleaned_text = ' '.join(page_text.split())
                            if len(cleaned_text) > 10:  # Only add substantial content
                                content_parts.append(f"[Page {page_num}] {cleaned_text}")
                    except Exception as page_error:
                        logger.warning(f"Error extracting page {page_num} from {file_path}: {str(page_error)}")
                        continue
                        
        except Exception as e:
            logger.error(f"Enhanced PDF extraction error for {file_path}: {str(e)}")
            # Try alternative extraction method
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages, 1):
                        text = page.extract_text()
                        if text and text.strip():
                            content_parts.append(f"[Page {page_num}] {text}")
            except ImportError:
                logger.warning("pdfplumber not available for fallback extraction")
            except Exception as fallback_error:
                logger.error(f"Fallback PDF extraction failed: {str(fallback_error)}")
        
        final_content = "\n\n".join(content_parts)
        
        # If still no content, log detailed diagnostic info
        if not final_content.strip():
            logger.error(f"No text content extracted from PDF: {file_path}")
            logger.error(f"File size: {os.path.getsize(file_path)} bytes")
            return f"[PDF Document: {os.path.basename(file_path)}] - Content extraction failed. This may be an image-based PDF or corrupted file."
        
        return final_content
    
    def _fallback_pdf_extraction(self, file_path: str) -> str:
        """Fallback PDF extraction method"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                content = ""
                for page in pdf_reader.pages:
                    content += page.extract_text() + "\n"
                return content.strip()
        except Exception as e:
            logger.error(f"Fallback PDF extraction failed for {file_path}: {str(e)}")
            return ""
    
    def _extract_docx_enhanced(self, file_path: str) -> str:
        """Enhanced DOCX extraction preserving structure"""
        try:
            doc = docx.Document(file_path)
            content_parts = []
            
            # Extract document properties
            if hasattr(doc.core_properties, 'title') and doc.core_properties.title:
                content_parts.append(f"Document Title: {doc.core_properties.title}")
            
            # Extract content with structure preservation
            current_section = []
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    if current_section:
                        content_parts.append("\n".join(current_section))
                        current_section = []
                    continue
                
                # Detect headers (simple heuristic)
                if (len(text) < 100 and 
                    (text.isupper() or 
                     paragraph.style.name.startswith('Heading') if hasattr(paragraph, 'style') else False)):
                    if current_section:
                        content_parts.append("\n".join(current_section))
                        current_section = []
                    content_parts.append(f"[Header] {text}")
                else:
                    current_section.append(text)
            
            # Add remaining content
            if current_section:
                content_parts.append("\n".join(current_section))
            
            return "\n\n".join(content_parts)
            
        except Exception as e:
            logger.error(f"Enhanced DOCX extraction error for {file_path}: {str(e)}")
            return self._fallback_docx_extraction(file_path)
    
    def _fallback_docx_extraction(self, file_path: str) -> str:
        """Fallback DOCX extraction"""
        try:
            doc = docx.Document(file_path)
            content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)
            return "\n".join(content)
        except Exception as e:
            logger.error(f"Fallback DOCX extraction failed for {file_path}: {str(e)}")
            return ""
    
    def _extract_text_enhanced(self, file_path: str) -> str:
        """Enhanced text extraction with encoding detection"""
        encodings_to_try = ['utf-8', 'utf-16', 'latin1', 'cp1252']
        
        for encoding in encodings_to_try:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    content = file.read()
                    # Add filename as context
                    filename = Path(file_path).stem
                    return f"[File: {filename}]\n\n{content}"
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"Text extraction error for {file_path}: {str(e)}")
                break
        
        return ""
    
    def _extract_markdown_enhanced(self, file_path: str) -> str:
        """Enhanced markdown extraction preserving structure"""
        content = self._extract_text_enhanced(file_path)
        
        if not content:
            return ""
        
        # Preserve markdown structure but make it more readable for GraphRAG
        # Convert headers to readable format
        content = re.sub(r'^#{1,6}\s+(.+)$', r'[Header] \1', content, flags=re.MULTILINE)
        
        # Preserve lists but make them cleaner
        content = re.sub(r'^\s*[-*+]\s+(.+)$', r'• \1', content, flags=re.MULTILINE)
        
        # Clean up multiple newlines
        content = re.sub(r'\n{3,}', '\n\n', content)
        
        return content
    
    def _extract_json_enhanced(self, file_path: str) -> str:
        """Enhanced JSON extraction with structure preservation"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
                
                # Create readable representation
                content_parts = [f"[JSON Document: {Path(file_path).name}]"]
                
                # Handle different JSON structures
                if isinstance(data, dict):
                    content_parts.append(self._format_json_dict(data))
                elif isinstance(data, list):
                    content_parts.append(f"JSON Array with {len(data)} items:")
                    for i, item in enumerate(data[:5]):  # Show first 5 items
                        content_parts.append(f"Item {i+1}: {str(item)[:200]}")
                else:
                    content_parts.append(f"JSON Value: {str(data)}")
                
                return "\n\n".join(content_parts)
                
        except Exception as e:
            logger.error(f"Enhanced JSON extraction error for {file_path}: {str(e)}")
            return self._fallback_json_extraction(file_path)
    
    def _format_json_dict(self, data: dict, max_depth: int = 3, current_depth: int = 0) -> str:
        """Format JSON dictionary for readability"""
        
        if current_depth >= max_depth:
            return f"[Complex object with {len(data)} keys]"
        
        parts = []
        for key, value in list(data.items())[:10]:  # Limit to first 10 keys
            if isinstance(value, dict):
                parts.append(f"{key}: {self._format_json_dict(value, max_depth, current_depth + 1)}")
            elif isinstance(value, list):
                parts.append(f"{key}: [Array with {len(value)} items]")
            else:
                parts.append(f"{key}: {str(value)[:100]}")
        
        return "\n".join(parts)
    
    def _fallback_json_extraction(self, file_path: str) -> str:
        """Fallback JSON extraction"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
                return json.dumps(data, indent=2)
        except Exception as e:
            logger.error(f"Fallback JSON extraction failed for {file_path}: {str(e)}")
            return ""
    
    def _extract_csv_enhanced(self, file_path: str) -> str:
        """Enhanced CSV extraction with structure analysis"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                lines = file.readlines()
                
            if not lines:
                return ""
            
            content_parts = [f"[CSV Document: {Path(file_path).name}]"]
            
            # Analyze CSV structure
            header_line = lines[0].strip()
            headers = [h.strip().strip('"') for h in header_line.split(',')]
            
            content_parts.append(f"CSV with {len(headers)} columns: {', '.join(headers)}")
            content_parts.append(f"Total rows: {len(lines) - 1}")
            
            # Show sample data
            if len(lines) > 1:
                content_parts.append("Sample data:")
                for i, line in enumerate(lines[1:6], 1):  # Show first 5 data rows
                    values = [v.strip().strip('"') for v in line.strip().split(',')]
                    row_summary = ', '.join(f"{h}: {v[:50]}" for h, v in zip(headers, values))
                    content_parts.append(f"Row {i}: {row_summary}")
            
            # Add raw content for detailed analysis
            content_parts.append("\n[Raw CSV Content]")
            content_parts.append(''.join(lines))
            
            return "\n\n".join(content_parts)
            
        except Exception as e:
            logger.error(f"Enhanced CSV extraction error for {file_path}: {str(e)}")
            return self._fallback_csv_extraction(file_path)
    
    def _fallback_csv_extraction(self, file_path: str) -> str:
        """Fallback CSV extraction"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            return ""
    
    def _extract_xlsx_enhanced(self, file_path: str) -> str:
        """Enhanced Excel extraction with sheet analysis"""
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True)
            content_parts = [f"[Excel Document: {Path(file_path).name}]"]
            content_parts.append(f"Sheets: {', '.join(wb.sheetnames)}")
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                content_parts.append(f"\n--- Sheet: {sheet_name} ---")
                
                rows_data = []
                for row in ws.iter_rows(max_row=min(50, ws.max_row), values_only=True):
                    if any(cell is not None for cell in row):
                        row_text = ', '.join([str(cell) if cell is not None else '' for cell in row])
                        rows_data.append(row_text)
                
                content_parts.extend(rows_data)
            
            return "\n".join(content_parts)
        except Exception:
            return f"[Excel Document: {Path(file_path).name}] - Content extraction failed"
    
    def _extract_pptx_enhanced(self, file_path: str) -> str:
        """Enhanced PowerPoint extraction with slide structure"""
        try:
            prs = Presentation(file_path)
            content_parts = [f"[PowerPoint Document: {Path(file_path).name}]"]
            content_parts.append(f"Total slides: {len(prs.slides)}")
            
            for i, slide in enumerate(prs.slides, 1):
                content_parts.append(f"\n--- Slide {i} ---")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text.strip())
                    
                    if shape.shape_type == 13 and hasattr(shape, 'table'):
                        table = shape.table
                        for row in table.rows:
                            row_text = ' | '.join([cell.text.strip() for cell in row.cells])
                            if row_text.strip():
                                content_parts.append(f"Table: {row_text}")
            
            return "\n".join(content_parts)
        except Exception:
            return f"[PowerPoint Document: {Path(file_path).name}] - Content extraction failed"
    
    def _extract_image_enhanced(self, file_path: str) -> str:
        """Enhanced image extraction with OCR"""
        try:
            img = Image.open(file_path)
            content_parts = [f"[Image Document: {Path(file_path).name}]"]
            content_parts.append(f"Image size: {img.size[0]}x{img.size[1]}")
            content_parts.append(f"Image mode: {img.mode}")
            
            try:
                text = pytesseract.image_to_string(img)
                if text.strip():
                    content_parts.append("Extracted text:")
                    content_parts.append(text.strip())
                else:
                    content_parts.append("No text detected in image")
            except Exception:
                content_parts.append("OCR extraction not available")
            
            return "\n".join(content_parts)
        except Exception:
            return f"[Image Document: {Path(file_path).name}] - Content extraction failed"
    
    def _classify_topic(self, content: str) -> str:
        """Classify document topic for GraphRAG organization"""
        content_lower = content.lower()
        
        # Enhanced topic classification
        topic_patterns = {
            'mathematical': ['theorem', 'proof', 'equation', 'formula', 'algorithm', 'mathematics'],
            'technical': ['system', 'architecture', 'software', 'code', 'implementation'],
            'scientific': ['research', 'study', 'experiment', 'hypothesis', 'analysis'],
            'business': ['business', 'market', 'strategy', 'revenue', 'customer'],
            'legal': ['law', 'legal', 'regulation', 'compliance', 'contract'],
            'medical': ['medical', 'health', 'treatment', 'diagnosis', 'patient'],
            'educational': ['course', 'lesson', 'tutorial', 'learning', 'education']
        }
        
        scores = {}
        for topic, keywords in topic_patterns.items():
            score = sum(1 for keyword in keywords if keyword in content_lower)
            if score > 0:
                scores[topic] = score
        
        return max(scores, key=scores.get) if scores else 'general'


# Alias for backwards compatibility
DocumentProcessor = EnhancedDocumentProcessor


class QueryMode(Enum):
    """Query execution modes"""
    READ = "READ"
    WRITE = "WRITE"
    SCHEMA = "SCHEMA"


class IndexType(Enum):
    """Types of database indexes"""
    BTREE = "BTREE"
    FULLTEXT = "FULLTEXT"
    VECTOR = "VECTOR"
    COMPOSITE = "COMPOSITE"


@dataclass
class DatabaseConfig:
    """Neo4j database configuration"""
    uri: str
    user: str
    password: str
    database: str = "neo4j"
    max_connection_lifetime: int = 3600
    max_connection_pool_size: int = 50
    connection_timeout: int = 30
    max_transaction_retry_time: int = 30
    encrypted: bool = False


@dataclass
class PerformanceConfig:
    """Performance tuning configuration"""
    batch_size: int = 1000
    index_creation_timeout: int = 300
    query_timeout: int = 30
    memory_limit: int = 4 * 1024 * 1024 * 1024  # 4GB
    cache_ttl: int = 3600  # 1 hour
    max_concurrent_operations: int = 10
    enable_query_logging: bool = True


@dataclass
class IsolationConfig:
    """Folder isolation configuration"""
    strict_folder_isolation: bool = True
    allow_cross_folder_read: bool = False
    validate_all_queries: bool = True
    audit_cross_folder_attempts: bool = True
    auto_fix_isolation_violations: bool = False


@dataclass
class FolderStatistics:
    """Enhanced statistics for GraphRAG folders"""
    folder_id: str
    folder_path: str
    entity_count: int
    relationship_count: int
    document_count: int
    community_count: int
    last_updated: datetime
    avg_entity_degree: float
    processing_sessions: int
    content_diversity_score: float
    graphrag_readiness_score: float


@dataclass
class QueryMetrics:
    """Enhanced query performance metrics"""
    query_hash: str
    execution_time: float
    records_returned: int
    records_available: int
    memory_usage: int
    index_hits: int
    search_mode: SearchMode
    communities_accessed: int
    entities_accessed: int
    folder_id: Optional[str] = None